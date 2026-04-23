from __future__ import annotations

from datetime import datetime
from pathlib import Path

from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Inches, Pt, RGBColor
from openpyxl import Workbook
from openpyxl.chart import BarChart, LineChart, PieChart, Reference
from openpyxl.formatting.rule import CellIsRule, ColorScaleRule
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.worksheet.datavalidation import DataValidation
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor as PPTXColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PInches, Pt as PPt
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY, TA_LEFT
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.platypus import (
    Image,
    KeepTogether,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)
from PIL import Image as PILImage

ROOT = Path(__file__).resolve().parents[1]
INPUTS = ROOT / "inputs"
ASSETS = ROOT / "assets"

PRIMARY = "173B68"
ACCENT = "2E86AB"
MINT = "5CC8B2"
LIGHT = "EDF4FA"
DARK = "22303F"


def ensure_dirs() -> None:
    INPUTS.mkdir(parents=True, exist_ok=True)
    ASSETS.mkdir(parents=True, exist_ok=True)


def copy_image_assets() -> tuple[Path, Path]:
    source_banner = Path("/mnt/data/a_clean_modern_abstract_business_analytics_graphi.png")
    source_isometric = Path("/mnt/data/a_clean_isometric_business_technology_infographic.png")
    banner = ASSETS / "analytics_banner.png"
    network = ASSETS / "logistics_isometric.png"
    if source_banner.exists():
        banner.write_bytes(source_banner.read_bytes())
    if source_isometric.exists():
        network.write_bytes(source_isometric.read_bytes())
    return banner, network


def add_hyperlink(paragraph, text: str, url: str, color: str = "0563C1"):
    part = paragraph.part
    r_id = part.relate_to(url, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink", is_external=True)
    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("r:id"), r_id)
    run = OxmlElement("w:r")
    r_pr = OxmlElement("w:rPr")
    color_el = OxmlElement("w:color")
    color_el.set(qn("w:val"), color)
    r_pr.append(color_el)
    underline = OxmlElement("w:u")
    underline.set(qn("w:val"), "single")
    r_pr.append(underline)
    run.append(r_pr)
    text_el = OxmlElement("w:t")
    text_el.text = text
    run.append(text_el)
    hyperlink.append(run)
    paragraph._p.append(hyperlink)
    return hyperlink


def set_cell_shading(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill)
    tc_pr.append(shd)


def style_doc(doc: Document) -> None:
    styles = doc.styles
    styles["Normal"].font.name = "Aptos"
    styles["Normal"].font.size = Pt(10.5)
    for name, size, color in [
        ("Title", 26, PRIMARY),
        ("Heading 1", 16, PRIMARY),
        ("Heading 2", 12, ACCENT),
        ("Subtitle", 12, DARK),
    ]:
        if name in styles:
            styles[name].font.name = "Aptos"
            styles[name].font.size = Pt(size)
            styles[name].font.color.rgb = RGBColor.from_string(color)


def create_docx(banner: Path, network: Path) -> Path:
    doc = Document()
    style_doc(doc)
    sec = doc.sections[0]
    sec.top_margin = Cm(1.8)
    sec.bottom_margin = Cm(1.5)
    sec.left_margin = Cm(2.0)
    sec.right_margin = Cm(2.0)
    header = sec.header.paragraphs[0]
    header.text = "Strategy & Transformation Office | Sample complex DOCX"
    header.runs[0].font.size = Pt(9)
    header.runs[0].font.color.rgb = RGBColor.from_string("6B7785")
    footer = sec.footer.paragraphs[0]
    footer.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    footer_run = footer.add_run("Confidential internal workshop pack")
    footer_run.font.size = Pt(9)
    footer_run.font.color.rgb = RGBColor.from_string("6B7785")

    p = doc.add_paragraph(style="Title")
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run("Q3 Transformation Program Brief")
    sp = doc.add_paragraph(style="Subtitle")
    sp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sp.add_run("Multi-section document with tables, images, bullets, quote block, hyperlink and page break")
    sp.runs[0].italic = True
    doc.add_paragraph("")
    if banner.exists():
        doc.add_picture(str(banner), width=Inches(6.2))
        doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER

    lead = doc.add_paragraph()
    lead.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    lead.add_run("Purpose. ").bold = True
    lead.add_run(
        "This sample document simulates a program brief for downstream Markdown conversion. "
        "It intentionally mixes narrative prose, a summary table, a pull-quote, numbered actions, "
        "an embedded illustration and a structured risk register so that Docling receives a rich DOCX input."
    )

    h1 = doc.add_paragraph(style="Heading 1")
    h1.add_run("1. Executive summary")
    summary_table = doc.add_table(rows=5, cols=3)
    summary_table.style = "Table Grid"
    widths = [Cm(4.0), Cm(5.5), Cm(5.5)]
    headers = ["Theme", "Current state", "Target by quarter-end"]
    for i, cell in enumerate(summary_table.rows[0].cells):
        cell.text = headers[i]
        set_cell_shading(cell, PRIMARY)
        for p in cell.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.color.rgb = RGBColor(255, 255, 255)
    rows = [
        ("Revenue enablement", "Pipeline quality improving but handoffs remain manual.", "Single operating cadence with stage-gated reviews and AI summaries."),
        ("Supplier resilience", "Critical vendors mapped; second-source coverage incomplete.", "Coverage above 85% for top spend categories."),
        ("Data governance", "Definitions differ across finance and operations dashboards.", "Unified KPI dictionary signed off by the PMO and finance controller."),
        ("Operating model", "Regional rituals exist but escalation paths are inconsistent.", "Regional playbook with measurable SLA and ownership matrix."),
    ]
    for ridx, row in enumerate(rows, start=1):
        for cidx, value in enumerate(row):
            summary_table.cell(ridx, cidx).text = value
            summary_table.cell(ridx, cidx).width = widths[cidx]
            if ridx % 2 == 0:
                set_cell_shading(summary_table.cell(ridx, cidx), "F7FAFD")
    doc.add_paragraph("")

    quote = doc.add_paragraph()
    quote.paragraph_format.left_indent = Cm(1.0)
    quote.paragraph_format.right_indent = Cm(0.5)
    quote.paragraph_format.space_before = Pt(6)
    quote.paragraph_format.space_after = Pt(8)
    run = quote.add_run(
        '"The operating model is now robust enough to scale, but only if we turn fragmented reporting into a disciplined weekly narrative."'
    )
    run.italic = True
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor.from_string(ACCENT)

    h2 = doc.add_paragraph(style="Heading 2")
    h2.add_run("Key workstreams")
    bullets = [
        "Standardise the executive review pack and align metric definitions across regions.",
        "Introduce a supplier watchlist combining risk, spend concentration and incident trend signals.",
        "Automate meeting notes and action extraction for recurring governance forums.",
    ]
    for item in bullets:
        doc.add_paragraph(item, style="List Bullet")

    num_items = [
        "Approve the revised KPI dictionary.",
        "Nominate accountable owners for each cross-functional dependency.",
        "Pilot document-to-Markdown conversion for steering packs and audit appendices.",
    ]
    for item in num_items:
        doc.add_paragraph(item, style="List Number")

    link_p = doc.add_paragraph()
    link_p.add_run("Reference repository: ")
    add_hyperlink(link_p, "Transformation playbook", "https://example.com/transformation-playbook")

    doc.add_paragraph(style="Heading 1").add_run("2. Delivery architecture")
    if network.exists():
        doc.add_picture(str(network), width=Inches(5.8))
        cap = doc.add_paragraph("Figure 1. Isometric network view used to test image extraction and caption retention.")
        cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
        cap.runs[0].italic = True
        cap.runs[0].font.size = Pt(9)
        cap.runs[0].font.color.rgb = RGBColor.from_string("5F6B78")

    doc.add_paragraph(
        "The architecture combines workstream-level plans, a common decision log and a shared dependency register. "
        "Each pillar publishes a Friday checkpoint note, while the PMO curates a Monday leadership digest with red-amber-green status and unresolved blockers."
    )

    doc.add_section(WD_SECTION.NEW_PAGE)
    sec2 = doc.sections[-1]
    sec2.top_margin = Cm(1.8)
    sec2.bottom_margin = Cm(1.5)
    sec2.left_margin = Cm(2.0)
    sec2.right_margin = Cm(2.0)
    doc.add_paragraph(style="Heading 1").add_run("3. Risk register")
    risk_table = doc.add_table(rows=6, cols=5)
    risk_table.style = "Table Grid"
    risk_headers = ["ID", "Risk", "Probability", "Impact", "Mitigation"]
    for i, text in enumerate(risk_headers):
        cell = risk_table.cell(0, i)
        cell.text = text
        set_cell_shading(cell, ACCENT)
        for p in cell.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.color.rgb = RGBColor(255, 255, 255)
    risks = [
        ("R-01", "Finance and operations disagree on pipeline qualification criteria.", "Medium", "High", "Ratify a single decision tree in the weekly PMO forum."),
        ("R-02", "Vendor remediation dates move without visible audit evidence.", "Medium", "Medium", "Add proof-of-completion attachment requirements."),
        ("R-03", "Legacy deck workflow causes version drift ahead of steering committee.", "High", "High", "Adopt source-controlled Markdown outputs and central approval."),
        ("R-04", "Regional teams duplicate analysis with different assumptions.", "Low", "Medium", "Publish one scenario model with approved inputs and refresh rhythm."),
        ("R-05", "Scanned appendices reduce extractability for downstream automation.", "Medium", "High", "Enable OCR only on image-based sections and preserve native text elsewhere."),
    ]
    for ridx, row in enumerate(risks, start=1):
        for cidx, value in enumerate(row):
            risk_table.cell(ridx, cidx).text = value
            if ridx % 2 == 1:
                set_cell_shading(risk_table.cell(ridx, cidx), "F9FBFC")

    doc.add_paragraph(style="Heading 2").add_run("Appendix note")
    doc.add_paragraph(
        "This final page intentionally mixes long table cells, headings and paragraph styles to help validate that the DOCX-to-Markdown path preserves hierarchy and table semantics."
    )

    path = INPUTS / "complex_program_brief.docx"
    doc.save(path)
    return path


def create_xlsx() -> Path:
    wb = Workbook()
    ws = wb.active
    ws.title = "Overview"
    ws.sheet_view.showGridLines = False

    title_fill = PatternFill("solid", fgColor=PRIMARY)
    section_fill = PatternFill("solid", fgColor=ACCENT)
    soft_fill = PatternFill("solid", fgColor=LIGHT)
    mint_fill = PatternFill("solid", fgColor="EAF8F4")
    white_font = Font(color="FFFFFF", bold=True, size=12)
    title_font = Font(color="FFFFFF", bold=True, size=16)
    section_font = Font(color="FFFFFF", bold=True, size=11)
    body_font = Font(color=DARK, size=10)
    thin = Side(style="thin", color="D6DEE8")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    ws.merge_cells("B2:H2")
    ws["B2"] = "Operations performance workbook"
    ws["B2"].fill = title_fill
    ws["B2"].font = title_font
    ws["B2"].alignment = Alignment(horizontal="center", vertical="center")
    ws.row_dimensions[2].height = 28

    ws.merge_cells("B3:H3")
    ws["B3"] = "Complex workbook with formulas, conditional formatting, validation, merged headers and charts"
    ws["B3"].fill = soft_fill
    ws["B3"].font = Font(color=PRIMARY, italic=True, size=10)
    ws["B3"].alignment = Alignment(horizontal="center")

    for col in range(2, 9):
        ws.column_dimensions[get_column_letter(col)].width = 16

    ws["B5"] = "KPI"
    ws["C5"] = "Current"
    ws["D5"] = "Target"
    ws["E5"] = "Variance"
    ws["F5"] = "Status"
    ws["G5"] = "Owner"
    ws["H5"] = "Commentary"
    for cell in ws["B5:H5"][0]:
        cell.fill = section_fill
        cell.font = white_font
        cell.alignment = Alignment(horizontal="center")
        cell.border = border

    kpis = [
        ("On-time shipment", 93, 96, "=C6-D6", "Amber", "Supply Chain", "Weather-related delays normalized in week 6."),
        ("Forecast accuracy", 88, 92, "=C7-D7", "Amber", "Finance", "Demand sensing pilot improving top account predictability."),
        ("Supplier coverage", 81, 85, "=C8-D8", "Red", "Procurement", "Second-source contracting remains open for packaging."),
        ("Order automation", 74, 80, "=C9-D9", "Green", "Operations", "New workflow lifted straight-through processing materially."),
    ]
    start_row = 6
    for r, row in enumerate(kpis, start=start_row):
        for c, value in enumerate(row, start=2):
            ws.cell(r, c, value)
            ws.cell(r, c).border = border
            ws.cell(r, c).font = body_font
            ws.cell(r, c).alignment = Alignment(vertical="center", wrap_text=True)
        if r % 2 == 0:
            for c in range(2, 9):
                ws.cell(r, c).fill = soft_fill
    for row in range(start_row, start_row + len(kpis)):
        ws.cell(row, 3).number_format = '0"%"'
        ws.cell(row, 4).number_format = '0"%"'
        ws.cell(row, 5).number_format = '0"%"'

    ws.conditional_formatting.add(f"F{start_row}:F{start_row+len(kpis)-1}", CellIsRule(operator='equal', formula=['"Green"'], fill=PatternFill("solid", fgColor="D9F2E6")))
    ws.conditional_formatting.add(f"F{start_row}:F{start_row+len(kpis)-1}", CellIsRule(operator='equal', formula=['"Amber"'], fill=PatternFill("solid", fgColor="FFF1CC")))
    ws.conditional_formatting.add(f"F{start_row}:F{start_row+len(kpis)-1}", CellIsRule(operator='equal', formula=['"Red"'], fill=PatternFill("solid", fgColor="F8D7DA")))

    ws["B12"] = "Weekly throughput"
    ws["B12"].fill = section_fill
    ws["B12"].font = section_font
    for i, hdr in enumerate(["Week", "Orders", "Backlog", "Cycle days"], start=2):
        cell = ws.cell(13, i, hdr)
        cell.fill = section_fill
        cell.font = white_font
        cell.border = border
        cell.alignment = Alignment(horizontal="center")
    throughput = [
        ("W1", 1210, 188, 5.9),
        ("W2", 1288, 176, 5.5),
        ("W3", 1345, 169, 5.1),
        ("W4", 1382, 160, 4.9),
        ("W5", 1410, 154, 4.7),
        ("W6", 1478, 143, 4.5),
    ]
    for ridx, row in enumerate(throughput, start=14):
        for cidx, value in enumerate(row, start=2):
            ws.cell(ridx, cidx, value)
            ws.cell(ridx, cidx).border = border
            ws.cell(ridx, cidx).font = body_font
            if ridx % 2 == 0:
                ws.cell(ridx, cidx).fill = mint_fill
    ws.conditional_formatting.add("B14:D19", ColorScaleRule(start_type='min', start_color='F8FBFD', mid_type='percentile', mid_value=50, mid_color='CDE7F7', end_type='max', end_color='7FB3D5'))

    line = LineChart()
    line.title = "Orders and backlog trend"
    line.height = 7
    line.width = 11
    data = Reference(ws, min_col=3, max_col=4, min_row=13, max_row=19)
    cats = Reference(ws, min_col=2, min_row=14, max_row=19)
    line.add_data(data, titles_from_data=True)
    line.set_categories(cats)
    line.y_axis.title = "Volume"
    line.x_axis.title = "Week"
    line.style = 2
    ws.add_chart(line, "J5")

    bar = BarChart()
    bar.title = "Cycle days"
    bar.height = 6
    bar.width = 7
    cycle = Reference(ws, min_col=5, max_col=5, min_row=13, max_row=19)
    bar.add_data(cycle, titles_from_data=True)
    bar.set_categories(cats)
    bar.style = 10
    ws.add_chart(bar, "J21")

    pipeline = wb.create_sheet("Pipeline")
    pipeline.sheet_view.showGridLines = False
    for col, width in zip(range(2, 9), [18, 14, 14, 16, 16, 16, 26]):
        pipeline.column_dimensions[get_column_letter(col)].width = width
    pipeline.merge_cells("B2:H2")
    pipeline["B2"] = "Initiative pipeline"
    pipeline["B2"].fill = title_fill
    pipeline["B2"].font = title_font
    pipeline["B2"].alignment = Alignment(horizontal="center")
    headers = ["Initiative", "Stage", "Owner", "Budget kEUR", "Benefit kEUR", "Confidence", "Notes"]
    for i, hdr in enumerate(headers, start=2):
        cell = pipeline.cell(4, i, hdr)
        cell.fill = section_fill
        cell.font = white_font
        cell.border = border
        cell.alignment = Alignment(horizontal="center")
    pipeline_rows = [
        ("Supplier watchtower", "Design", "Procurement", 180, 420, "Medium", "Data model approved; pilot scope in review."),
        ("Order automation", "Build", "Operations", 240, 670, "High", "E2E workflow validated in plant A."),
        ("Narrative reporting", "Pilot", "PMO", 95, 220, "High", "Markdown outputs replace manual steering notes."),
        ("Demand sensing", "Ideation", "Commercial", 150, 360, "Low", "External data feeds pending legal review."),
        ("Returns analytics", "Scale", "Finance", 110, 280, "Medium", "Control group reached statistical significance."),
    ]
    for ridx, row in enumerate(pipeline_rows, start=5):
        for cidx, value in enumerate(row, start=2):
            pipeline.cell(ridx, cidx, value)
            pipeline.cell(ridx, cidx).border = border
            pipeline.cell(ridx, cidx).font = body_font
            pipeline.cell(ridx, cidx).alignment = Alignment(vertical="center", wrap_text=True)
            if ridx % 2 == 1:
                pipeline.cell(ridx, cidx).fill = soft_fill
    dv = DataValidation(type="list", formula1='"Ideation,Design,Build,Pilot,Scale"', allow_blank=False)
    pipeline.add_data_validation(dv)
    dv.add("C5:C9")
    conf_dv = DataValidation(type="list", formula1='"Low,Medium,High"', allow_blank=False)
    pipeline.add_data_validation(conf_dv)
    conf_dv.add("G5:G9")

    summary = wb.create_sheet("Scenario")
    summary.sheet_view.showGridLines = False
    summary.merge_cells("B2:F2")
    summary["B2"] = "Scenario model"
    summary["B2"].fill = title_fill
    summary["B2"].font = title_font
    summary["B2"].alignment = Alignment(horizontal="center")
    headers = ["Assumption", "Base", "Stretch", "Downside", "Formula check"]
    for i, hdr in enumerate(headers, start=2):
        cell = summary.cell(4, i, hdr)
        cell.fill = section_fill
        cell.font = white_font
        cell.border = border
    assumptions = [
        ("Revenue uplift %", 3.5, 5.0, 1.2, "=C5-D5"),
        ("Inventory days", 41, 36, 46, "=C6-D6"),
        ("Working capital", 12.8, 14.2, 10.9, "=C7-D7"),
    ]
    for ridx, row in enumerate(assumptions, start=5):
        for cidx, value in enumerate(row, start=2):
            summary.cell(ridx, cidx, value)
            summary.cell(ridx, cidx).border = border
            summary.cell(ridx, cidx).font = body_font
            if ridx % 2 == 0:
                summary.cell(ridx, cidx).fill = soft_fill
    pie = PieChart()
    pie.title = "Investment mix"
    labels = Reference(summary, min_col=2, min_row=5, max_row=7)
    pie_data = Reference(summary, min_col=3, min_row=4, max_row=7)
    pie.add_data(pie_data, titles_from_data=True)
    pie.set_categories(labels)
    pie.height = 6
    pie.width = 7
    ws_ref = summary
    ws_ref.add_chart(pie, "H4")

    path = INPUTS / "complex_operations_dashboard.xlsx"
    wb.save(path)
    return path


def create_pptx(banner: Path, network: Path) -> Path:
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)

    def add_title(slide, title: str, subtitle: str | None = None):
        tx = slide.shapes.add_textbox(PInches(0.6), PInches(0.35), PInches(7.5), PInches(0.8))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = PPt(28)
        p.font.bold = True
        p.font.color.rgb = PPTXColor.from_string(PRIMARY)
        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.size = PPt(12)
            p2.font.color.rgb = PPTXColor.from_string("5B6875")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if banner.exists():
        slide.shapes.add_picture(str(banner), 0, 0, width=prs.slide_width, height=prs.slide_height)
    overlay = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, PInches(0.55), PInches(0.55), PInches(5.4), PInches(2.1))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = PPTXColor.from_string("FFFFFF")
    overlay.fill.transparency = 10
    overlay.line.color.rgb = PPTXColor.from_string("D7E2EC")
    add_title(slide, "Executive Steering Committee Update", "Sample PPTX crafted to stress Markdown extraction from slides")
    box = slide.shapes.add_textbox(PInches(0.8), PInches(1.75), PInches(4.9), PInches(1.2))
    tf = box.text_frame
    for idx, text in enumerate([
        "• Operating rhythm redesigned around weekly narrative packs",
        "• Supplier risk watchlist launched in pilot markets",
        "• Source material includes charts, images, tables and callouts",
    ]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = PPt(14)
        p.font.color.rgb = PPTXColor.from_string(DARK)
        p.level = 0
    footer = slide.shapes.add_textbox(PInches(0.8), PInches(6.65), PInches(3.3), PInches(0.3))
    footer.text_frame.paragraphs[0].text = "Confidential | Strategy demo deck"
    footer.text_frame.paragraphs[0].font.size = PPt(9)
    footer.text_frame.paragraphs[0].font.color.rgb = PPTXColor.from_string("5B6875")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Performance snapshot", "The slide mixes native charts, KPI cards and explanatory callouts")
    cards = [
        (0.7, "On-time shipment", "93%", "Target 96%"),
        (3.1, "Forecast accuracy", "88%", "Improved +4pp"),
        (5.5, "Supplier coverage", "81%", "Gap in packaging"),
    ]
    for x, label, value, note in cards:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PInches(x), PInches(1.35), PInches(2.1), PInches(1.35))
        shape.fill.solid()
        shape.fill.fore_color.rgb = PPTXColor.from_string("F4F8FB")
        shape.line.color.rgb = PPTXColor.from_string("D6E1EA")
        tf = shape.text_frame
        p1 = tf.paragraphs[0]
        p1.text = label
        p1.font.size = PPt(12)
        p1.font.color.rgb = PPTXColor.from_string(PRIMARY)
        p2 = tf.add_paragraph()
        p2.text = value
        p2.font.size = PPt(24)
        p2.font.bold = True
        p2.font.color.rgb = PPTXColor.from_string(ACCENT)
        p3 = tf.add_paragraph()
        p3.text = note
        p3.font.size = PPt(10)
        p3.font.color.rgb = PPTXColor.from_string("5B6875")
    chart_data = CategoryChartData()
    chart_data.categories = ["W1", "W2", "W3", "W4", "W5", "W6"]
    chart_data.add_series("Orders", (1210, 1288, 1345, 1382, 1410, 1478))
    chart_data.add_series("Backlog", (188, 176, 169, 160, 154, 143))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, PInches(0.8), PInches(3.0), PInches(6.1), PInches(3.3), chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.value_axis.has_title = True
    chart.value_axis.axis_title.text_frame.text = "Volume"
    chart.category_axis.tick_labels.font.size = PPt(10)
    note_shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PInches(7.3), PInches(1.55), PInches(5.1), PInches(4.6))
    note_shape.fill.solid()
    note_shape.fill.fore_color.rgb = PPTXColor.from_string("EDF7FF")
    note_shape.line.color.rgb = PPTXColor.from_string("B7D3EA")
    tf = note_shape.text_frame
    tf.word_wrap = True
    items = [
        "Signal: throughput improved six consecutive weeks.",
        "Implication: backlog is reducing without service quality erosion.",
        "Action: convert the analysis narrative into Markdown for the steering memo.",
    ]
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = PPt(15 if idx == 0 else 13)
        p.font.color.rgb = PPTXColor.from_string(DARK)
        if idx == 0:
            p.font.bold = True
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Architecture and control points", "Mixed layout with image, matrix table and milestone band")
    if network.exists():
        slide.shapes.add_picture(str(network), PInches(0.7), PInches(1.3), width=PInches(5.4), height=PInches(4.2))
    table = slide.shapes.add_table(4, 3, PInches(6.45), PInches(1.55), PInches(5.9), PInches(2.3)).table
    headers = ["Layer", "Owner", "Control objective"]
    for c, hdr in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = PPTXColor.from_string(PRIMARY)
        for p in cell.text_frame.paragraphs:
            p.font.size = PPt(12)
            p.font.bold = True
            p.font.color.rgb = PPTXColor.from_string("FFFFFF")
    table_rows = [
        ("Intake", "PMO", "Ensure source files are versioned and indexed."),
        ("Conversion", "Automation", "Preserve heading hierarchy, tables and captions."),
        ("Review", "Leadership", "Approve narrative and actions before distribution."),
    ]
    for r, row in enumerate(table_rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = PPTXColor.from_string("F7FAFD" if r % 2 else "EDF4FA")
            for p in cell.text_frame.paragraphs:
                p.font.size = PPt(11)
                p.font.color.rgb = PPTXColor.from_string(DARK)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, PInches(6.45), PInches(4.35), PInches(5.9), PInches(0.5))
    band.fill.solid()
    band.fill.fore_color.rgb = PPTXColor.from_string(ACCENT)
    band.line.color.rgb = PPTXColor.from_string(ACCENT)
    tf = band.text_frame
    p = tf.paragraphs[0]
    p.text = "Milestones | Scope freeze -> Pilot -> Governance rollout -> KPI audit"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = PPt(13)
    p.font.color.rgb = PPTXColor.from_string("FFFFFF")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Risks and next steps", "Intentional combination of bullets, action owners and closing statement")
    left = slide.shapes.add_textbox(PInches(0.8), PInches(1.5), PInches(6.0), PInches(4.6))
    tf = left.text_frame
    risk_items = [
        "Risk 1: scanned appendices may require OCR fallback.",
        "Risk 2: legacy deck content could flatten semantic hierarchy.",
        "Risk 3: regional variants may duplicate the same narrative with minor edits.",
        "Action: standardise one source package and convert each artifact to Markdown automatically.",
    ]
    for idx, item in enumerate(risk_items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = PPt(16 if idx == 0 else 14)
        p.font.color.rgb = PPTXColor.from_string(DARK)
        p.level = 0
    owner_table = slide.shapes.add_table(4, 2, PInches(7.1), PInches(1.6), PInches(5.1), PInches(2.9)).table
    owner_table.cell(0, 0).text = "Owner"
    owner_table.cell(0, 1).text = "Deliverable"
    for c in range(2):
        cell = owner_table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PPTXColor.from_string(PRIMARY)
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.color.rgb = PPTXColor.from_string("FFFFFF")
    owner_rows = [
        ("PMO lead", "Finalize source inventory and naming rules"),
        ("Automation lead", "Run Docling conversion and quality checks"),
        ("Executive sponsor", "Approve pilot scope and review cadence"),
    ]
    for r, row in enumerate(owner_rows, start=1):
        for c, value in enumerate(row):
            cell = owner_table.cell(r, c)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = PPTXColor.from_string("F7FAFD" if r % 2 else "EDF4FA")
    close = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PInches(7.1), PInches(5.0), PInches(5.1), PInches(1.05))
    close.fill.solid()
    close.fill.fore_color.rgb = PPTXColor.from_string("EAF8F4")
    close.line.color.rgb = PPTXColor.from_string(MINT)
    close.text_frame.paragraphs[0].text = "Desired test outcome: four distinct source formats -> four clean Markdown outputs in one target directory."
    close.text_frame.paragraphs[0].font.size = PPt(14)
    close.text_frame.paragraphs[0].font.color.rgb = PPTXColor.from_string(PRIMARY)

    path = INPUTS / "complex_steering_committee_update.pptx"
    prs.save(path)
    return path


def pdf_footer(canvas, doc):
    canvas.saveState()
    canvas.setFont("Helvetica", 9)
    canvas.setFillColor(colors.HexColor("#6B7785"))
    canvas.drawString(2 * cm, 1 * cm, "Supplier resilience assessment | Sample complex PDF")
    canvas.drawRightString(A4[0] - 2 * cm, 1 * cm, f"Page {doc.page}")
    canvas.restoreState()


def create_pdf(banner: Path, network: Path) -> Path:
    path = INPUTS / "complex_supplier_resilience_assessment.pdf"
    doc = SimpleDocTemplate(str(path), pagesize=A4, topMargin=1.8 * cm, bottomMargin=1.5 * cm, leftMargin=2.0 * cm, rightMargin=2.0 * cm)
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name="TitleCustom", parent=styles["Title"], fontName="Helvetica-Bold", fontSize=22, textColor=colors.HexColor(f"#{PRIMARY}"), alignment=TA_CENTER, spaceAfter=10))
    styles.add(ParagraphStyle(name="SubCustom", parent=styles["BodyText"], fontName="Helvetica", fontSize=11, textColor=colors.HexColor("#5B6875"), alignment=TA_CENTER, spaceAfter=8))
    styles.add(ParagraphStyle(name="BodyJustify", parent=styles["BodyText"], fontName="Helvetica", fontSize=10.3, leading=14, alignment=TA_JUSTIFY, textColor=colors.HexColor(f"#{DARK}")))
    styles.add(ParagraphStyle(name="Heading1Custom", parent=styles["Heading1"], fontName="Helvetica-Bold", fontSize=15, textColor=colors.HexColor(f"#{PRIMARY}"), spaceBefore=10, spaceAfter=8))
    styles.add(ParagraphStyle(name="Heading2Custom", parent=styles["Heading2"], fontName="Helvetica-Bold", fontSize=11.5, textColor=colors.HexColor(f"#{ACCENT}"), spaceBefore=8, spaceAfter=6))
    styles.add(ParagraphStyle(name="Callout", parent=styles["BodyText"], fontName="Helvetica-Oblique", fontSize=10.5, backColor=colors.HexColor("#EDF7FF"), textColor=colors.HexColor(f"#{ACCENT}"), leftIndent=10, rightIndent=10, borderPadding=8, borderWidth=0.5, borderColor=colors.HexColor("#B7D3EA"), leading=13))

    story = []
    story.append(Paragraph("Supplier Resilience Assessment", styles["TitleCustom"]))
    story.append(Paragraph("Complex PDF input with banner image, tables, multi-page narrative and figure caption", styles["SubCustom"]))
    if banner.exists():
        story.append(Image(str(banner), width=16.0 * cm, height=8.2 * cm))
        story.append(Spacer(1, 0.25 * cm))
    story.append(Paragraph(
        "This PDF is deliberately rich in structure so it can serve as a realistic source file for Docling. It includes a cover section, a narrative summary, a formatted table, a second-page illustration and several callout blocks designed to test Markdown heading retention and reading order.",
        styles["BodyJustify"],
    ))
    story.append(Spacer(1, 0.2 * cm))
    story.append(Paragraph("1. Assessment summary", styles["Heading1Custom"]))
    story.append(Paragraph(
        "The review focused on supplier concentration, route complexity, packaging dependency and incident recovery speed. Across the top twelve spend categories, single-source exposure remains material in specialized packaging, while transport resilience improved following the introduction of regional contingency lanes.",
        styles["BodyJustify"],
    ))
    story.append(Spacer(1, 0.15 * cm))
    story.append(Paragraph(
        '"Resilience does not come from more reporting alone; it comes from consistent evidence, clear thresholds and fast escalation when assumptions break."',
        styles["Callout"],
    ))
    story.append(Spacer(1, 0.2 * cm))
    story.append(Paragraph("2. Heatmap table", styles["Heading1Custom"]))

    data = [
        ["Category", "Exposure", "Lead time", "Recovery readiness", "Commentary"],
        ["Packaging", "High", "24 days", "Medium", "Dual-source contracting still in negotiation."],
        ["Electronics", "Medium", "18 days", "High", "Component alternates qualified in two regions."],
        ["Logistics", "Medium", "11 days", "Medium", "Seasonal capacity remains sensitive to weather."],
        ["Labeling", "Low", "7 days", "High", "High standardization enables quick rerouting."],
    ]
    tbl = Table(data, colWidths=[3.1 * cm, 2.3 * cm, 2.5 * cm, 3.0 * cm, 6.0 * cm])
    tbl.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(f"#{PRIMARY}")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, -1), 9.2),
        ("BACKGROUND", (0, 1), (-1, -1), colors.HexColor("#F7FAFD")),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.HexColor("#F7FAFD"), colors.HexColor("#EDF4FA")]),
        ("GRID", (0, 0), (-1, -1), 0.35, colors.HexColor("#D6DEE8")),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("LEADING", (0, 0), (-1, -1), 11),
    ]))
    story.append(tbl)
    story.append(Spacer(1, 0.25 * cm))
    story.append(Paragraph("Priority actions", styles["Heading2Custom"]))
    for item in [
        "Approve packaging second-source business case.",
        "Add documentary proof to all remediation closures.",
        "Convert appendix source files into Markdown and version them centrally.",
    ]:
        story.append(Paragraph(f"• {item}", styles["BodyJustify"]))
    story.append(PageBreak())
    story.append(Paragraph("3. Network view and reading-order stress test", styles["Heading1Custom"]))
    if network.exists():
        story.append(Image(str(network), width=15.4 * cm, height=8.3 * cm))
        story.append(Paragraph("Figure 1. Network-style logistics illustration included to test image references in extracted Markdown.", styles["SubCustom"]))
    story.append(Paragraph(
        "On the second page, the content intentionally alternates between image, caption, heading and dense narrative text. This helps evaluate whether a downstream converter preserves logical order rather than flattening the page into a visually convenient but semantically wrong sequence.",
        styles["BodyJustify"],
    ))
    story.append(Spacer(1, 0.18 * cm))
    story.append(Paragraph(
        "A robust converter should recognize headings, paragraph groups, bullet-like statements and captions. For born-digital PDFs, OCR should generally remain off so the system can rely on the native text layer, fonts and layout coordinates instead of re-reading the page as an image.",
        styles["BodyJustify"],
    ))
    story.append(Spacer(1, 0.15 * cm))
    story.append(Paragraph(
        "Closing note: this document is intended purely as a complex test asset and does not contain real operational data.",
        styles["Callout"],
    ))
    doc.build(story, onFirstPage=pdf_footer, onLaterPages=pdf_footer)
    return path


def main() -> None:
    ensure_dirs()
    banner, network = copy_image_assets()
    outputs = [
        create_docx(banner, network),
        create_xlsx(),
        create_pptx(banner, network),
        create_pdf(banner, network),
    ]
    print("Created sample files:")
    for path in outputs:
        print(path)


if __name__ == "__main__":
    main()
